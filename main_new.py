import os
import re
import fitz  # PyMuPDF
import docx
import openpyxl
from pptx import Presentation
import cv2
import pytesseract
import numpy as np

# --- Folders ---
input_folder = r"E:\project_optive\input_files"
redacted_folder = r"E:\project_optive\redacted_files"
os.makedirs(redacted_folder, exist_ok=True)

# --- Sensitive Patterns ---
SENSITIVE_PATTERNS = {
    # Contact Information
    "EMAIL": r"\b[a-zA-Z0-9_.+-]+@[a-zA-Z0-9-]+\.[a-zA-Z0-9-.]+\b",
    "PHONE": r"\(?\d{3}\)?[-.\s]?\d{3}[-.\s]?\d{4}",
    "PHONE_INTL": r"\+\d{1,3}[-.\s]?\(?\d{1,4}\)?[-.\s]?\d{1,4}[-.\s]?\d{1,9}",
    # Government IDs
    "AADHAAR": r"\b\d{4}\s?\d{4}\s?\d{4}\b",
    "PAN": r"\b[A-Z]{5}[0-9]{4}[A-Z]\b",
    "SSN": r"\b\d{3}-\d{2}-\d{4}\b",
    "PASSPORT": r"\b[A-Z]{1,2}\d{6,9}\b",
    # Financial
    "CREDIT_CARD": r"\b(?:\d{4}[-\s]?){3}\d{4}\b",
    "BANK_ACCOUNT": r"\b\d{9,18}\b",
    "IFSC": r"\b[A-Z]{4}0[A-Z0-9]{6}\b",
    "SWIFT": r"\b[A-Z]{6}[A-Z0-9]{2}([A-Z0-9]{3})?\b",
    # Network & Security
    "IPV4": r"\b(?:(?:25[0-5]|2[0-4][0-9]|[01]?[0-9][0-9]?)\.){3}(?:25[0-5]|2[0-4][0-9]|[01]?[0-9][0-9]?)(?:/\d{1,2})?\b",
    "IPV6": r"\b(?:[0-9a-fA-F]{1,4}:){7}[0-9a-fA-F]{1,4}\b",
    "MAC_ADDRESS": r"\b(?:[0-9A-Fa-f]{2}[:-]){5}(?:[0-9A-Fa-f]{2})\b",
    "PORT": r":\d{2,5}\b(?=\s|,|$|\)|})",
    # AWS & Cloud
    "AWS_ACCESS_KEY": r"\b(?:AKIA|ASIA|AIDA|AROA|AIPA|ANPA|ANVA|APKA)[A-Z0-9]{16}\b",
    "AWS_SECRET_KEY": r"\b[A-Za-z0-9/+=]{40}\b(?=\s|,|\"|}|$)",
    "AWS_ARN": r"\barn:aws:[a-z0-9\-]+:[a-z0-9\-]*:\d{12}:[a-zA-Z0-9\-/:]+",
    "AWS_ACCOUNT_ID": r"\b\d{12}(?=:|\"|\s|,|}|\))",
    "GOOGLE_API_KEY": r"\bAIza[0-9A-Za-z\-_]{35}\b",
    # Tokens & Keys
    "JWT_TOKEN": r"\beyJ[A-Za-z0-9_\-]+\.eyJ[A-Za-z0-9_\-]+\.[A-Za-z0-9_\-]+",
    "BEARER_TOKEN": r"\bBearer\s+[A-Za-z0-9\-._~+/]+=*",
    "GITHUB_TOKEN": r"\bghp_[A-Za-z0-9]{36}\b",
    "GITHUB_OAUTH": r"\bgho_[A-Za-z0-9]{36}\b",
    "API_KEY": r"\b[A-Za-z0-9_\-]{32,45}\b",
    "PRIVATE_KEY": r"-----BEGIN (?:RSA |EC |OPENSSH )?PRIVATE KEY-----[^-]+-----END (?:RSA |EC |OPENSSH )?PRIVATE KEY-----",
    # Database & Connection Strings
    "DB_CONNECTION": r"(?i)(?:mongodb|mysql|postgresql|jdbc|redis):\/\/[^\s\"'<>]+",
    "PASSWORD_FIELD": r"(?i)(?:password|passwd|pwd|secret)[\"\s]*[:=][\"\s]*[^\s\"}{,]{6,}",
    # Dates & IDs
    "DOB": r"\b\d{1,2}[/-]\d{1,2}[/-]\d{2,4}\b",
    "DATE_ISO": r"\b\d{4}-\d{2}-\d{2}(?:T\d{2}:\d{2}:\d{2}(?:\.\d+)?(?:Z|[+-]\d{2}:\d{2}))?\b",
    "VEHICLE": r"\b[A-Z]{2}[ -]?\d{1,2}[ -]?[A-Z]{1,2}[ -]?\d{4}\b",
    # Names & IDs...
    # Add rest of your patterns here as needed
}

# --- Redact Text ---
def redact_text(text):
    for pattern in SENSITIVE_PATTERNS.values():
        text = re.sub(pattern, "[REDACTED]", text, flags=re.IGNORECASE)
    return text

# --- Summarize non-sensitive content ---
def summarize_text(text):
    return " ".join([w for w in text.split() if w != "[REDACTED]"])

# --- File Type Redaction Functions ---
def redact_pdf(file_path, save_path):
    doc = fitz.open(file_path)
    summary = ""
    for page in doc:
        text = page.get_text()
        redacted = redact_text(text)
        page.clean_contents()
        page.insert_text((50,50), redacted, fontsize=11)
        summary += summarize_text(redacted) + " "
    doc.save(save_path)
    return summary.strip()

def redact_docx(file_path, save_path):
    doc = docx.Document(file_path)
    summary = ""
    for para in doc.paragraphs:
        redacted = redact_text(para.text)
        para.text = redacted
        summary += summarize_text(redacted) + " "
    doc.save(save_path)
    return summary.strip()

def redact_xlsx(file_path, save_path):
    wb = openpyxl.load_workbook(file_path)
    summary = ""
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        for row in ws.iter_rows():
            for cell in row:
                if cell.value:
                    redacted = redact_text(str(cell.value))
                    cell.value = redacted
                    summary += summarize_text(redacted) + " "
    wb.save(save_path)
    return summary.strip()

def redact_pptx(file_path, save_path):
    prs = Presentation(file_path)
    summary = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                redacted = redact_text(shape.text)
                shape.text = redacted
                summary += summarize_text(redacted) + " "
    prs.save(save_path)
    return summary.strip()

def redact_image(file_path, save_path):
    img = cv2.imread(file_path)
    if img is None:
        return ""
    gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
    data = pytesseract.image_to_data(gray, output_type=pytesseract.Output.DICT)
    summary = ""
    for i in range(len(data['text'])):
        word = data['text'][i]
        if word:
            redacted_word = redact_text(word)
            if "[REDACTED]" in redacted_word:
                x, y, w, h = data['left'][i], data['top'][i], data['width'][i], data['height'][i]
                cv2.rectangle(img, (x, y), (x + w, y + h), (0,0,0), -1)
            summary += summarize_text(redacted_word) + " "
    # Blur faces
    face_cascade = cv2.CascadeClassifier(cv2.data.haarcascades + "haarcascade_frontalface_default.xml")
    faces = face_cascade.detectMultiScale(gray, 1.1, 4)
    for (x, y, w, h) in faces:
        img[y:y+h, x:x+w] = cv2.GaussianBlur(img[y:y+h, x:x+w], (51,51), 30)
    cv2.imwrite(save_path, img)
    return summary.strip()

# --- Process Files ---
for file_name in os.listdir(input_folder):
    file_path = os.path.join(input_folder, file_name)
    redacted_path = os.path.join(redacted_folder, file_name)
    try:
        if file_name.lower().endswith(".pdf"):
            summary = redact_pdf(file_path, redacted_path)
        elif file_name.lower().endswith(".docx"):
            summary = redact_docx(file_path, redacted_path)
        elif file_name.lower().endswith(".xlsx"):
            summary = redact_xlsx(file_path, redacted_path)
        elif file_name.lower().endswith(".pptx"):
            summary = redact_pptx(file_path, redacted_path)
        elif file_name.lower().endswith((".png",".jpg",".jpeg")):
            summary = redact_image(file_path, redacted_path)
        else:
            print(f"{file_name}: Unsupported format")
            continue
        print(f"{file_name}: Redacted and saved.")
        print(f"Summary of non-sensitive content: {summary}\n")
    except Exception as e:
        print(f"Error processing {file_name}: {e}")
