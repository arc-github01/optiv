import os
import io
import re
import fitz
import docx
from pathlib import Path
import pandas as pd
import cv2
import numpy as np
from PIL import Image, ImageFilter
import pytesseract
from facenet_pytorch import MTCNN
import spacy
import openpyxl
from openpyxl.styles import PatternFill, Font
from pptx import Presentation
import gradio as gr
import google.generativeai as genai

pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

nlp = spacy.load("en_core_web_sm")
mtcnn = MTCNN(keep_all=True, device="cpu")
genai.configure(api_key="your_api_key")


# Comprehensive sensitive patterns
SENSITIVE_PATTERNS = {
    # Contact Information
    "EMAIL": r"\b[a-zA-Z0-9_.+-]+@[a-zA-Z0-9-]+\.[a-zA-Z0-9-.]+\b",
    "PHONE": r"\(?\d{3}\)?[-.\s]?\d{3}[-.\s]?\d{4}",
    "PHONE_INTL": r"\+\d{1,3}[-.\s]?\(?\d{1,4}\)?[-.\s]?\d{1,4}[-.\s]?\d{1,9}",
    
    # Government IDs
    "AADHAAR": r"\b\d{4}\s?\d{4}\s?\d{4}\b",
    "PAN": r"\b[A-Z]{5}[0-9]{4}[A-Z]\b",
    "SSN": r"\b\d{3}-\d{2}-\d{4}\b",
    "PASSPORT": r"\b[A-Z]{1,2}\d{6,9}\b",
    
    # Financial
    "CREDIT_CARD": r"\b(?:\d{4}[-\s]?){3}\d{4}\b",
    "BANK_ACCOUNT": r"\b\d{9,18}\b",
    "IFSC": r"\b[A-Z]{4}0[A-Z0-9]{6}\b",
    "SWIFT": r"\b[A-Z]{6}[A-Z0-9]{2}([A-Z0-9]{3})?\b",
    
    # Network & Security
    "IPV4": r"\b(?:(?:25[0-5]|2[0-4][0-9]|[01]?[0-9][0-9]?)\.){3}(?:25[0-5]|2[0-4][0-9]|[01]?[0-9][0-9]?)(?:/\d{1,2})?\b",
    "IPV6": r"\b(?:[0-9a-fA-F]{1,4}:){7}[0-9a-fA-F]{1,4}\b",
    "MAC_ADDRESS": r"\b(?:[0-9A-Fa-f]{2}[:-]){5}(?:[0-9A-Fa-f]{2})\b",
    "PORT": r":\d{2,5}\b(?=\s|,|$|\)|})",
    
    # AWS & Cloud
    "AWS_ACCESS_KEY": r"\b(?:AKIA|ASIA|AIDA|AROA|AIPA|ANPA|ANVA|APKA)[A-Z0-9]{16}\b",
    "AWS_SECRET_KEY": r"\b[A-Za-z0-9/+=]{40}\b(?=\s|,|\"|}|$)",
    "AWS_ARN": r"\barn:aws:[a-z0-9\-]+:[a-z0-9\-]*:\d{12}:[a-zA-Z0-9\-/:]+",
    "AWS_ACCOUNT_ID": r"\b\d{12}(?=:|\"|\s|,|}|\))",
    "GOOGLE_API_KEY": r"\bAIza[0-9A-Za-z\-_]{35}\b",
    
    # Tokens & Keys
    "JWT_TOKEN": r"\beyJ[A-Za-z0-9_\-]+\.eyJ[A-Za-z0-9_\-]+\.[A-Za-z0-9_\-]+",
    "BEARER_TOKEN": r"\bBearer\s+[A-Za-z0-9\-._~+/]+=*",
    "GITHUB_TOKEN": r"\bghp_[A-Za-z0-9]{36}\b",
    "GITHUB_OAUTH": r"\bgho_[A-Za-z0-9]{36}\b",
    "API_KEY": r"\b[A-Za-z0-9_\-]{32,45}\b",
    "PRIVATE_KEY": r"-----BEGIN (?:RSA |EC |OPENSSH )?PRIVATE KEY-----[^-]+-----END (?:RSA |EC |OPENSSH )?PRIVATE KEY-----",
    
    # Database & Connection Strings
    "DB_CONNECTION": r"(?i)(?:mongodb|mysql|postgresql|jdbc|redis):\/\/[^\s\"'<>]+",
    "PASSWORD_FIELD": r"(?i)(?:password|passwd|pwd|secret)[\"\s]*[:=][\"\s]*[^\s\"}{,]{6,}",
    
    # Dates & IDs
    "DOB": r"\b\d{1,2}[/-]\d{1,2}[/-]\d{2,4}\b",
    "DATE_ISO": r"\b\d{4}-\d{2}-\d{2}(?:T\d{2}:\d{2}:\d{2}(?:\.\d+)?(?:Z|[+-]\d{2}:\d{2}))?\b",
    "VEHICLE": r"\b[A-Z]{2}[ -]?\d{1,2}[ -]?[A-Z]{1,2}[ -]?\d{4}\b",
    
    # Postal Codes
    "ZIPCODE_US": r"\b\d{5}(?:-\d{4})?\b",
    "POSTAL_CODE_UK": r"\b[A-Z]{1,2}\d{1,2}[A-Z]?\s?\d[A-Z]{2}\b",
    "PINCODE_INDIA": r"\b\d{6}\b",
    "POSTAL_CODE_CANADA": r"\b[A-Z]\d[A-Z]\s?\d[A-Z]\d\b",
    
    # Addresses
    "ADDRESS_STREET": r"\b\d+\s+(?:[A-Z][a-z]+\s+){1,5}(?:Street|St|Avenue|Ave|Road|Rd|Boulevard|Blvd|Lane|Ln|Drive|Dr|Court|Ct|Circle|Cir|Way|Place|Pl|Parkway|Pkwy|Highway|Hwy|Terrace|Ter)\b",
    "ADDRESS_PO_BOX": r"\bP\.?O\.?\s*Box\s+\d+\b",
    "ADDRESS_UNIT": r"\b(?:Apt|Apartment|Suite|Unit|#)\s*[A-Z0-9\-]+\b",
    "ADDRESS_BUILDING": r"\b(?:Building|Block|Tower|Wing)\s+[A-Z0-9\-]+\b",
    "ADDRESS_FLOOR": r"\b(?:Floor|Flr|Level)\s+\d+\b",
    "ADDRESS_CITY_STATE": r"\b[A-Z][a-z]+(?:\s+[A-Z][a-z]+)*,\s*[A-Z]{2}\s+\d{5}\b",
    "ADDRESS_FULL_US": r"\d+\s+[A-Za-z0-9\s,]+,\s*[A-Z][a-z]+,\s*[A-Z]{2}\s+\d{5}",
    "ADDRESS_INDIAN": r"\b(?:House|H\.?No\.?|Plot|Shop)\s*[#:]?\s*[\d\-/A-Z]+[,\s]+(?:[A-Z][a-z]+\s*){1,4}(?:Nagar|Colony|Road|Street|Area|Layout|Extension|Ext|Society|Sector)",
    
    # Names & Personal Info
    "NAME_AFTER_KEYWORD": r"(?i)(?:name|naam|नाम|student\s*name|employee\s*name|full\s*name)[\s:=-]+([A-Za-zा-ॿ\s]{2,50})",
    "FATHER_NAME": r"(?i)(?:father|पिता|father'?s?\s*name)[\s:=-]+([A-Za-zा-ॿ\s]{2,50})",
    "MOTHER_NAME": r"(?i)(?:mother|माता|mother'?s?\s*name)[\s:=-]+([A-Za-zा-ॿ\s]{2,50})",
    
    # Student/Employee IDs
    "STUDENT_ID": r"(?i)(?:student\s*id|student\s*no|roll\s*no|roll\s*number|enrollment\s*no)[\s:=-]+([A-Za-z0-9\-/]{3,20})",
    "EMPLOYEE_ID": r"(?i)(?:employee\s*id|emp\s*id|staff\s*id|employee\s*no|emp\s*no)[\s:=-]+([A-Za-z0-9\-/]{3,20})",
    "GENERIC_ID": r"(?i)(?:id\s*number|id\s*no|identification)[\s:=-]+([A-Za-z0-9\-/]{3,20})",
    "ALPHANUMERIC_ID_1": r"\b(?:EMP|STU|STUD|ROLL|ENR|ID)[A-Z]*\d{4,10}\b",
    "ALPHANUMERIC_ID_2": r"\b\d{2,4}[A-Z]{2,6}\d{4,10}\b",
    "ALPHANUMERIC_ID_3": r"\b[A-Z]{2,4}\d{6,10}\b",
    "ALPHANUMERIC_ID_4": r"\b\d{4,6}[A-Z]{1,3}\d{3,6}\b",
    
    # JSON/Structured Data
    "JSON_ROLE_ID": r'"RoleId"\s*:\s*"[A-Z0-9]+"',
    "JSON_ROLE_NAME": r'"RoleName"\s*:\s*"[^"]+"',
    "JSON_ARN": r'"Arn"\s*:\s*"arn:aws:[^"]+"',
    "JSON_SID": r'"Sid"\s*:\s*"[^"]+"',
    "JSON_AWS": r'"AWS"\s*:\s*"arn:aws:[^"]+"',
    "JSON_ACCOUNT": r'"AccountId"\s*:\s*"\d{12}"',
    "JSON_KEY_VALUE": r'"[A-Za-z]{4,}"\s*:\s*"[^"]{10,}"',
}


def extract_from_pdf(path):
    try:
        doc = fitz.open(path)
        text = ""
        for page in doc:
            text += page.get_text()
        doc.close()
        return text
    except Exception as e:
        print(f"Error reading PDF: {e}")
        import traceback
        traceback.print_exc()
        return ""


def extract_from_docx(path):
    docx_file = docx.Document(path)
    text_content = []
    
    for paragraph in docx_file.paragraphs:
        text_content.append(paragraph.text)
    
    for table in docx_file.tables:
        for row in table.rows:
            row_text = [cell.text for cell in row.cells]
            text_content.append("\t".join(row_text))
    
    return "\n".join(text_content)


def extract_from_excel(path):
    try:
        if path.endswith('.xlsx'):
            wb = openpyxl.load_workbook(path, data_only=True)
            text_content = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_content.append(f"\n=== Sheet: {sheet_name} ===\n")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = [str(cell) if cell is not None else "" for cell in row]
                    if any(row_text):
                        text_content.append("\t".join(row_text))
            
            return "\n".join(text_content)
        else:
            excel_file = pd.ExcelFile(path)
            text_content = []
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(path, sheet_name=sheet_name)
                text_content.append(f"\n=== Sheet: {sheet_name} ===\n")
                df = df.fillna('')
                text_content.append(df.to_string(index=False))
            return "\n".join(text_content)
    except Exception as e:
        print(f"Error reading Excel file: {e}")
        import traceback
        traceback.print_exc()
        return f"Error reading Excel file: {e}"


def extract_from_ppt(path):
    try:
        prs = Presentation(path)
        text_content = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            text_content.append(f"\n=== Slide {slide_num} ===\n")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_content.append(shape.text)
                
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = [cell.text for cell in row.cells]
                        text_content.append("\t".join(row_text))
                
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text:
                                text_content.append(run.text)
        
        return "\n".join(text_content)
    except Exception as e:
        print(f"Error reading PowerPoint file: {e}")
        import traceback
        traceback.print_exc()
        return f"Error reading PowerPoint file: {e}"


def extract_from_image(path):
    img = Image.open(path).convert("RGB")
    text = pytesseract.image_to_string(img)
    return text, img


def is_sensitive_text(text):
    """Check if text contains sensitive information"""
    if not text or len(str(text)) < 2:
        return False
    
    text_str = str(text)
    
    for pattern_name, pattern in SENSITIVE_PATTERNS.items():
        try:
            if re.search(pattern, text_str, re.IGNORECASE):
                return True
        except:
            pass
    
    try:
        if len(text_str) > 2 and not text_str.isdigit():
            doc = nlp(text_str)
            if any(ent.label_ in ["PERSON", "ORG", "GPE", "LOC"] for ent in doc.ents):
                return True
    except:
        pass
    
    return False


def redact_text(text):
    """Redact sensitive information from text"""
    if not text:
        return text
    
    priority_patterns = [
        "PRIVATE_KEY", "JWT_TOKEN", "BEARER_TOKEN", "AWS_ARN", 
        "AWS_ACCESS_KEY", "AWS_SECRET_KEY", "GOOGLE_API_KEY",
        "GITHUB_TOKEN", "GITHUB_OAUTH", "DB_CONNECTION", "PASSWORD_FIELD",
        "JSON_ARN", "JSON_ROLE_ID", "JSON_ROLE_NAME", "JSON_SID", "JSON_AWS", "JSON_ACCOUNT",
        "EMAIL", "IPV4", "IPV6", "MAC_ADDRESS", 
        "AWS_ACCOUNT_ID", "AADHAAR", "PAN", "CREDIT_CARD", "SSN",
        "PHONE", "PHONE_INTL", "VEHICLE", "DATE_ISO", "DOB", 
        "ADDRESS_FULL_US", "ADDRESS_INDIAN", "ADDRESS_STREET", "ADDRESS_CITY_STATE",
        "ADDRESS_PO_BOX", "ADDRESS_UNIT", "ADDRESS_BUILDING", "ADDRESS_FLOOR",
        "ZIPCODE_US", "POSTAL_CODE_UK", "PINCODE_INDIA", "POSTAL_CODE_CANADA",
        "NAME_AFTER_KEYWORD", "FATHER_NAME", "MOTHER_NAME",
        "STUDENT_ID", "EMPLOYEE_ID", "GENERIC_ID",
        "ALPHANUMERIC_ID_1", "ALPHANUMERIC_ID_2", "ALPHANUMERIC_ID_3", "ALPHANUMERIC_ID_4"
    ]
    
    for label in priority_patterns:
        if label in SENSITIVE_PATTERNS:
            text = re.sub(SENSITIVE_PATTERNS[label], f"[REDACTED_{label}]", text, flags=re.IGNORECASE)
    
    for label, pattern in SENSITIVE_PATTERNS.items():
        if label not in priority_patterns:
            text = re.sub(pattern, f"[REDACTED_{label}]", text, flags=re.IGNORECASE)
    
    try:
        doc = nlp(text)
        redacted_positions = []
        for ent in doc.ents:
            if ent.label_ in ["PERSON", "GPE", "ORG", "LOC", "FAC"]:
                overlap = any(not (ent.start_char >= end or ent.end_char <= start) for start, end in redacted_positions)
                if not overlap and len(ent.text) > 2:
                    text = text[:ent.start_char] + f"[REDACTED_{ent.label_}]" + text[ent.end_char:]
                    new_end = ent.start_char + len(f"[REDACTED_{ent.label_}]")
                    redacted_positions.append((ent.start_char, new_end))
    except Exception as e:
        print(f"SpaCy NER error: {e}")
    
    return text


def redact_excel(input_path, output_path):
    try:
        if input_path.endswith('.xlsx'):
            wb = openpyxl.load_workbook(input_path)
            black_fill = PatternFill(start_color="000000", end_color="000000", fill_type="solid")
            white_font = Font(color="FFFFFF")
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value is None:
                            continue
                        
                        if is_sensitive_text(str(cell.value)):
                            cell.fill = black_fill
                            cell.font = white_font
                            cell.value = "[REDACTED]"
            
            wb.save(output_path)
            return output_path
        else:
            output_path = output_path.replace('.xls', '.xlsx')
            excel_file = pd.ExcelFile(input_path)
            with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
                for sheet_name in excel_file.sheet_names:
                    df = pd.read_excel(input_path, sheet_name=sheet_name)
                    
                    for col in df.columns:
                        df[col] = df[col].apply(lambda x: "[REDACTED]" if pd.notna(x) and is_sensitive_text(str(x)) else x)
                    
                    df.to_excel(writer, sheet_name=sheet_name, index=False)
            
            return output_path
            
    except Exception as e:
        print(f"Error redacting Excel: {e}")
        import traceback
        traceback.print_exc()
        return None


def redact_pdf(input_path, output_path):
    """Create redacted PDF by overlaying black rectangles on sensitive text and blurring images"""
    try:
        doc = fitz.open(input_path)
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            
            # Redact text
            text_instances = page.get_text("dict")
            blocks = text_instances.get("blocks", [])
            for block in blocks:
                if "lines" in block:
                    for line in block["lines"]:
                        line_text = " ".join([span.get("text", "") for span in line["spans"]])
                        if is_sensitive_text(line_text):
                            for span in line["spans"]:
                                rect = fitz.Rect(span["bbox"])
                                page.draw_rect(rect, color=(0, 0, 0), fill=(0, 0, 0))
                        else:
                            for span in line["spans"]:
                                text = span.get("text", "")
                                if is_sensitive_text(text):
                                    rect = fitz.Rect(span["bbox"])
                                    page.draw_rect(rect, color=(0, 0, 0), fill=(0, 0, 0))
            
            # Process images with better error handling
            image_list = page.get_images(full=True)
            for img_index, img_info in enumerate(image_list):
                try:
                    xref = img_info[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    
                    pil_image = Image.open(io.BytesIO(image_bytes))
                    processed_image = pil_image.convert("RGB")
                    
                    # Blur small images
                    if pil_image.width < 500 and pil_image.height < 500:
                        processed_image = processed_image.filter(ImageFilter.GaussianBlur(20))
                    
                    # Apply face detection and code detection
                    processed_image = blur_faces(processed_image)
                    processed_image = detect_and_blur_codes(processed_image)
                    
                    # Save processed image
                    img_byte_arr = io.BytesIO()
                    processed_image.save(img_byte_arr, format='PNG')
                    img_byte_arr.seek(0)
                    
                    # Try to get image rectangle, with fallback
                    try:
                        img_rect = page.get_image_rects(img_info)[0]
                    except (IndexError, Exception) as e:
                        # Fallback: try to find image position using alternative method
                        print(f"Could not get image rect for image {img_index}, trying alternative method")
                        
                        # Try to get image bbox using xref
                        try:
                            img_rect = page.get_image_bbox(img_info)
                        except:
                            # If both methods fail, skip this image
                            print(f"Skipping image {img_index} - could not determine position")
                            continue
                    
                    # Validate rectangle before inserting
                    if img_rect and img_rect.is_valid and not img_rect.is_empty:
                        page.insert_image(img_rect, stream=img_byte_arr.getvalue(), keep_proportion=False)
                    else:
                        print(f"Invalid image rectangle for image {img_index}, skipping")
                        
                except Exception as img_error:
                    print(f"Error processing image {img_index} on page {page_num}: {img_error}")
                    continue
        
        doc.save(output_path)
        doc.close()
        return output_path
        
    except Exception as e:
        print(f"Error redacting PDF: {e}")
        import traceback
        traceback.print_exc()
        return None


def redact_docx(input_path, output_path):
    """Create redacted DOCX by replacing sensitive text"""
    try:
        doc = docx.Document(input_path)
        
        for paragraph in doc.paragraphs:
            if paragraph.text:
                redacted_text = redact_text(paragraph.text)
                if redacted_text != paragraph.text:
                    paragraph.text = redacted_text
        
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text:
                        redacted_text = redact_text(cell.text)
                        if redacted_text != cell.text:
                            cell.text = redacted_text
        
        doc.save(output_path)
        return output_path
    except Exception as e:
        print(f"Error redacting DOCX: {e}")
        import traceback
        traceback.print_exc()
        return None


def redact_ppt(input_path, output_path):
    try:
        prs = Presentation(input_path)
        
        for slide in prs.slides:
            shapes_to_replace = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    original_text = shape.text
                    redacted_text = redact_text(original_text)
                    
                    if shape.has_text_frame and original_text != redacted_text:
                        shape.text_frame.clear()
                        p = shape.text_frame.paragraphs[0] if shape.text_frame.paragraphs else shape.text_frame.add_paragraph()
                        run = p.add_run()
                        run.text = redacted_text
                
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            if cell.text:
                                original_text = cell.text
                                redacted_text = redact_text(original_text)
                                if original_text != redacted_text:
                                    cell.text_frame.clear()
                                    p = cell.text_frame.paragraphs[0] if cell.text_frame.paragraphs else cell.text_frame.add_paragraph()
                                    run = p.add_run()
                                    run.text = redacted_text
                
                if shape.shape_type == 13:
                    shapes_to_replace.append(shape)

            for shape in shapes_to_replace:
                image = shape.image
                image_bytes = image.blob
                pil_image = Image.open(io.BytesIO(image_bytes))
                processed_image = pil_image.convert("RGB")
                
                if pil_image.width < 500 and pil_image.height < 500:
                    processed_image = processed_image.filter(ImageFilter.GaussianBlur(20))
                
                processed_image = blur_faces(processed_image)
                processed_image = detect_and_blur_codes(processed_image)
                
                with io.BytesIO() as output:
                    processed_image.save(output, format='PNG')
                    slide.shapes.add_picture(output, shape.left, shape.top, width=shape.width, height=shape.height)
                    sp = shape._element
                    sp.getparent().remove(sp)

        prs.save(output_path)
        return output_path
    except Exception as e:
        print(f"Error redacting PowerPoint: {e}")
        import traceback
        traceback.print_exc()
        return None

def redact_image_text(img):
    """Redact sensitive text in images with enhanced multi-word pattern matching"""
    
    # Ensure img is a valid PIL Image
    if not isinstance(img, Image.Image):
        if isinstance(img, np.ndarray):
            img = Image.fromarray(img)
        else:
            raise TypeError(f"Expected PIL Image or numpy array, got {type(img)}")
    
    # Convert to RGB if needed
    if img.mode != 'RGB':
        img = img.convert('RGB')
    
    full_text = pytesseract.image_to_string(img)
    
    is_aws_json = bool(re.search(r'"(?:RoleId|Arn|AWS|AccountId|Sid|Effect|Principal|Action)"', full_text, re.IGNORECASE))
    
    data = pytesseract.image_to_data(img, output_type=pytesseract.Output.DICT)
    
    blocks = {}
    for i in range(len(data['level'])):
        if data['text'][i].strip():
            block_num = data['block_num'][i]
            if block_num not in blocks:
                blocks[block_num] = []
            blocks[block_num].append(i)
    
    redacted_regions = []
    
    # AWS JSON special handling
    if is_aws_json:
        for i in range(len(data['text'])):
            text = data['text'][i].strip()
            if text and len(text) > 2:
                if re.search(r'(?:AKIA|ASIA|arn:aws|RoleId|Statement|Principal|Allow|Deny|\d{12}|"[A-Za-z]+")', text, re.IGNORECASE):
                    x, y, w, h = data['left'][i], data['top'][i], data['width'][i], data['height'][i]
                    padding = 8
                    x = max(0, x - padding)
                    y = max(0, y - padding)
                    w = min(img.width - x, w + 2 * padding)
                    h = min(img.height - y, h + 2 * padding)
                    
                    if w > 0 and h > 0:
                        region = img.crop((x, y, x + w, y + h))
                        region = region.filter(ImageFilter.GaussianBlur(30))
                        img.paste(region, (x, y))
                        redacted_regions.append((x, y, w, h))
    
    # Pattern-based redaction
    for block_num, word_indices in blocks.items():
        block_text = " ".join([data['text'][i] for i in word_indices])
        
        for pattern_name, pattern in SENSITIVE_PATTERNS.items():
            try:
                for match in re.finditer(pattern, block_text, re.IGNORECASE):
                    start, end = match.span()
                    
                    current_pos = 0
                    matched_indices = []
                    for i in word_indices:
                        word = data['text'][i]
                        word_len = len(word)
                        if start <= current_pos < end or start < current_pos + word_len <= end:
                            matched_indices.append(i)
                        current_pos += word_len + 1
                    
                    if matched_indices:
                        x1 = min(data['left'][i] for i in matched_indices)
                        y1 = min(data['top'][i] for i in matched_indices)
                        x2 = max(data['left'][i] + data['width'][i] for i in matched_indices)
                        y2 = max(data['top'][i] + data['height'][i] for i in matched_indices)
                        
                        padding = 10
                        x = max(0, x1 - padding)
                        y = max(0, y1 - padding)
                        w = min(img.width - x, (x2 - x1) + 2 * padding)
                        h = min(img.height - y, (y2 - y1) + 2 * padding)
                        
                        overlaps = False
                        for rx, ry, rw, rh in redacted_regions:
                            if not (x + w < rx or x > rx + rw or y + h < ry or y > ry + rh):
                                overlaps = True
                                break
                        
                        if not overlaps and w > 0 and h > 0:
                            region = img.crop((x, y, x + w, y + h))
                            region = region.filter(ImageFilter.GaussianBlur(30))
                            img.paste(region, (x, y))
                            redacted_regions.append((x, y, w, h))
            except Exception:
                continue
    
    # NER-based redaction
    for block_num, word_indices in blocks.items():
        block_text = " ".join([data['text'][i] for i in word_indices])
        
        try:
            doc = nlp(block_text)
            for ent in doc.ents:
                if ent.label_ in ["PERSON", "ORG"] and len(ent.text) > 2:
                    ent_start = block_text.find(ent.text)
                    if ent_start != -1:
                        ent_end = ent_start + len(ent.text)
                        
                        current_pos = 0
                        matched_indices = []
                        for i in word_indices:
                            word = data['text'][i]
                            word_len = len(word)
                            if ent_start <= current_pos < ent_end or ent_start < current_pos + word_len <= ent_end:
                                matched_indices.append(i)
                            current_pos += word_len + 1
                        
                        if matched_indices:
                            x1 = min(data['left'][i] for i in matched_indices)
                            y1 = min(data['top'][i] for i in matched_indices)
                            x2 = max(data['left'][i] + data['width'][i] for i in matched_indices)
                            y2 = max(data['top'][i] + data['height'][i] for i in matched_indices)
                            
                            padding = 10
                            x = max(0, x1 - padding)
                            y = max(0, y1 - padding)
                            w = min(img.width - x, (x2 - x1) + 2 * padding)
                            h = min(img.height - y, (y2 - y1) + 2 * padding)
                            
                            overlaps = False
                            for rx, ry, rw, rh in redacted_regions:
                                if not (x + w < rx or x > rx + rw or y + h < ry or y > ry + rh):
                                    overlaps = True
                                    break
                            
                            if not overlaps and w > 0 and h > 0:
                                region = img.crop((x, y, x + w, y + h))
                                region = region.filter(ImageFilter.GaussianBlur(30))
                                img.paste(region, (x, y))
                                redacted_regions.append((x, y, w, h))
        except Exception:
            continue
    
    return img


def detect_and_blur_codes(img):
    """QR/Barcode detection disabled - returns image unchanged"""
    return img

def blur_faces(img):
    """Detect and blur faces in images - with error handling"""
    try:
        boxes, _ = mtcnn.detect(img)
        if boxes is not None and len(boxes) > 0:  # Check if boxes is not empty
            for box in boxes:
                x1, y1, x2, y2 = map(int, box)
                x1, y1 = max(0, x1), max(0, y1)
                x2, y2 = min(img.width, x2), min(img.height, y2)
                
                if x2 > x1 and y2 > y1:
                    face = img.crop((x1, y1, x2, y2))
                    face = face.filter(ImageFilter.GaussianBlur(30))
                    img.paste(face, (x1, y1))
    except Exception as e:
        print(f"Face detection error: {e}")
    
    return img


def analyze_with_gemini(text, file_type, file_name):
    """Analyze document content using Gemini AI"""
    try:
        model = genai.GenerativeModel("gemini-2.0-flash-exp")
        
        prompt = f"""Analyze the following content from a {file_type} file named "{file_name}".

Content:
{text}

Please provide a structured analysis in the following format:
1. File Description: A brief description of what this file contains (2-3 sentences)
2. Key Findings: List 3-4 bullet points of the most important findings or insights

Be concise and focus on the main purpose and key information."""

        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        return f"Analysis failed: {e}"


def describe_image_with_gemini(img_path):
    """Analyze image content using Gemini AI"""
    try:
        model = genai.GenerativeModel("gemini-2.0-flash-exp")
        with open(img_path, "rb") as f:
            img_bytes = f.read()

        prompt = """Analyze this image and provide:
1. File Description: What is shown in the image? (2-3 sentences)
2. Key Findings: List 3-4 key observations or insights about what's in the image

Format your response clearly with these two sections."""

        response = model.generate_content([
            prompt,
            {"mime_type": "image/png", "data": img_bytes}
        ])
        return response.text
    except Exception as e:
        return f"Image analysis failed: {e}"


def process_single_file(file_path):
    """Process a single file: extract text, redact sensitive info, analyze content"""
    fname = str(file_path)
    file_name = os.path.basename(fname)
    ext = fname.split(".")[-1].lower()
    
    result = {
        "File Name": file_name,
        "File Type": f".{ext}",
        "File Description": "",
        "Key Findings": "",
        "Processed Image": None,
        "Redacted File": None
    }

    try:
        output_dir = os.path.join(os.path.dirname(fname), "redacted_output")
        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(output_dir, f"redacted_{file_name}")
        
        if ext == "pdf":
            raw_text = extract_from_pdf(fname)
            if not raw_text or len(raw_text.strip()) < 10:
                result["File Description"] = "PDF appears empty or could not be read"
                result["Key Findings"] = "Check if PDF is valid and not password protected"
                return result
                
            redacted_text = redact_text(raw_text)
            analysis = analyze_with_gemini(redacted_text[:5000], "PDF", file_name)
            
            redacted_file = redact_pdf(fname, output_path)
            result["Redacted File"] = redacted_file
            
        elif ext == "docx":
            raw_text = extract_from_docx(fname)
            redacted_text = redact_text(raw_text)
            analysis = analyze_with_gemini(redacted_text[:5000], "DOCX", file_name)
            
            redacted_file = redact_docx(fname, output_path)
            result["Redacted File"] = redacted_file
            
        elif ext in ["xlsx", "xls"]:
            raw_text = extract_from_excel(fname)
            redacted_text = redact_text(raw_text)
            analysis = analyze_with_gemini(redacted_text[:5000], "Excel", file_name)
            
            redacted_file = redact_excel(fname, output_path)
            result["Redacted File"] = redacted_file
            
        elif ext in ["pptx", "ppt"]:
            raw_text = extract_from_ppt(fname)
            redacted_text = redact_text(raw_text)
            analysis = analyze_with_gemini(redacted_text[:5000], "PowerPoint", file_name)
            
            redacted_file = redact_ppt(fname, output_path)
            result["Redacted File"] = redacted_file
            
        elif ext in ["jpg", "jpeg", "png"]:
            ocr_text, img = extract_from_image(fname)
            img = detect_and_blur_codes(img)
            img = redact_image_text(img)
            img = blur_faces(img)
            result["Processed Image"] = img
            
            img.save(output_path)
            result["Redacted File"] = output_path
            
            if len(ocr_text.strip()) < 10:
                analysis = describe_image_with_gemini(fname)
            else:
                redacted_text = redact_text(ocr_text)
                analysis = analyze_with_gemini(redacted_text[:3000], "Image", file_name)
        else:
            result["File Description"] = "Unsupported file type"
            return result
        
        lines = analysis.split('\n')
        description_lines = []
        findings_lines = []
        in_description = False
        in_findings = False
        
        for line in lines:
            line = line.strip()
            if 'file description' in line.lower() or line.startswith('1.'):
                in_description = True
                in_findings = False
                if ':' in line:
                    desc_text = line.split(':', 1)[1].strip()
                    if desc_text:
                        description_lines.append(desc_text)
                continue
            elif 'key findings' in line.lower() or line.startswith('2.'):
                in_description = False
                in_findings = True
                continue
            
            if in_description and line:
                description_lines.append(line)
            elif in_findings and line:
                findings_lines.append(line)
        
        result["File Description"] = ' '.join(description_lines) if description_lines else "Analysis completed"
        result["Key Findings"] = '\n'.join(findings_lines) if findings_lines else "See description"
        
    except Exception as e:
        result["File Description"] = f"Error processing file"
        result["Key Findings"] = str(e)
        import traceback
        traceback.print_exc()
    
    return result


def process_directory(directory_path):
    """Process all supported files in a directory"""
    if not directory_path or not os.path.exists(directory_path):
        return "Invalid directory path", None, None
    
    supported_extensions = ['.pdf', '.docx', '.xlsx', '.xls', '.pptx', '.ppt', '.jpg', '.jpeg', '.png']
    files = []
    
    path_obj = Path(directory_path)
    for ext in supported_extensions:
        files.extend(path_obj.glob(f'*{ext}'))
        files.extend(path_obj.glob(f'*{ext.upper()}'))
    
    if not files:
        return "No supported files found in directory", None, None
    
    results = []
    processed_images = []
    redacted_files_list = []
    
    for file_path in sorted(files):
        print(f"Processing: {file_path}")
        result = process_single_file(file_path)
        results.append(result)
        if result["Processed Image"] is not None:
            processed_images.append((result["File Name"], result["Processed Image"]))
        if result["Redacted File"] is not None:
            redacted_files_list.append(result["Redacted File"])
    
    df = pd.DataFrame(results)
    df = df[["File Name", "File Type", "File Description", "Key Findings"]]
    
    html_output = """
    <style>
        .analysis-table {
            width: 100%;
            border-collapse: collapse;
            font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
            margin-top: 20px;
            box-shadow: 0 8px 32px rgba(0,0,0,0.3);
            border-radius: 10px;
            overflow: hidden;
        }
        .analysis-table th {
            ackground: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            color: white;
            font-weight: 600;
            padding: 15px;
            text-align: left;
            border: none;
            font-size: 1.1em;
        }
       .analysis-table td {
            padding: 15px;
            border: none;
            vertical-align: top;
            line-height: 1.4;
            color: #e0e0e0;
        }
        .analysis-table tr:nth-child(even) {
            background-color: #2d3748;
        }
        .analysis-table tr:nth-child(odd) {
            background-color: #1a202c;
        }
        .analysis-table tr:hover {
            background-color: #4a5568;
            transition: background-color 0.3s ease;
        }
        
        .file-type {
            font-weight: bold;
            color: #63b3ed;
        }
        .key-findings {
            white-space: pre-line;
            color: #cbd5e0;
        }
        .key-findings::before {
            content: "• ";
            color: #68d391;
        }
    </style>
    
    <table class="analysis-table">
        <thead>
            <tr>
                <td style="color: #ffffff; font-weight: 600;">{row['File Name']}</td>
                <td class="file-type">{row['File Type']}</td>
                <td style="color: #e2e8f0;">{row['File Description']}</td>
                <td class="key-findings">{findings_formatted}</td>
            </tr>
        </thead>
        <tbody>
    """
    
    for _, row in df.iterrows():
        findings_formatted = row['Key Findings'].replace('\n', '\n• ')
        if findings_formatted and not findings_formatted.startswith('•'):
            findings_formatted = '• ' + findings_formatted
        
        html_output += f"""
            <tr>
                <td>{row['File Name']}</td>
                <td class="file-type">{row['File Type']}</td>
                <td>{row['File Description']}</td>
                <td class="key-findings">{findings_formatted}</td>
            </tr>
        """
    
    html_output += """
        </tbody>
    </table>
    """
    
    image_gallery = None
    if processed_images:
        image_gallery = [img for _, img in processed_images]
    
    download_info = f"Redacted files saved in: {os.path.join(directory_path, 'redacted_output')}" if redacted_files_list else "No files required redaction"
    
    return html_output, image_gallery, download_info

custom_css = """
:root {
    --primary: #4361ee;
    --secondary: #3a0ca3;
    --accent: #7209b7;
    --success: #4cc9f0;
    --warning: #f72585;
    --dark-bg: #0f0f23;
    --dark-card: #1a1a2e;
    --dark-text: #e2e8f0;
    --dark-border: #2d3748;
}

.gradio-container {
    font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
    background: linear-gradient(135deg, #0f0f23 0%, #1a1a2e 50%, #16213e 100%) !important;
    min-height: 100vh;
    color: var(--dark-text);
}

.main-container {
    background: rgba(26, 26, 46, 0.95) !important;
    backdrop-filter: blur(10px);
    border-radius: 20px;
    box-shadow: 0 20px 40px rgba(0,0,0,0.3);
    margin: 20px;
    padding: 30px;
    border: 1px solid var(--dark-border);
}

.header {
    text-align: center;
    margin-bottom: 30px;
    padding: 30px;
    background: linear-gradient(135deg, var(--primary), var(--secondary));
    border-radius: 15px;
    color: white;
    box-shadow: 0 10px 20px rgba(67, 97, 238, 0.3);
    border: 1px solid rgba(255,255,255,0.1);
}

.header h1 {
    margin: 0;
    font-size: 2.5em;
    font-weight: 700;
    text-shadow: 2px 2px 4px rgba(0,0,0,0.3);
}

.header p {
    margin: 10px 0 0 0;
    font-size: 1.2em;
    opacity: 0.9;
}

.feature-grid {
    display: grid;
    grid-template-columns: repeat(auto-fit, minmax(250px, 1fr));
    gap: 20px;
    margin: 30px 0;
}

.feature-card {
    background: var(--dark-card);
    padding: 25px;
    border-radius: 15px;
    text-align: center;
    box-shadow: 0 8px 16px rgba(0,0,0,0.2);
    transition: transform 0.3s ease, box-shadow 0.3s ease;
    border-left: 5px solid var(--primary);
    border: 1px solid var(--dark-border);
    color: var(--dark-text);
}

.feature-card:hover {
    transform: translateY(-5px);
    box-shadow: 0 15px 30px rgba(0,0,0,0.4);
    border-left: 5px solid var(--accent);
}

.feature-icon {
    font-size: 2.5em;
    margin-bottom: 15px;
    color: var(--primary);
}

.feature-card h3 {
    margin: 0 0 10px 0;
    color: white;
    font-weight: 600;
}

.feature-card p {
    margin: 0;
    color: #cbd5e0;
    line-height: 1.5;
}

.input-section {
    background: var(--dark-card);
    padding: 30px;
    border-radius: 15px;
    box-shadow: 0 8px 16px rgba(0,0,0,0.2);
    margin: 20px 0;
    border: 1px solid var(--dark-border);
}

.input-section h2 {
    color: white;
    margin-bottom: 20px;
    font-weight: 600;
}

.directory-input {
    background: #2d3748;
    border: 2px dashed #4a5568;
    border-radius: 10px;
    padding: 20px;
    transition: all 0.3s ease;
    color: white;
}

.directory-input:hover {
    border-color: var(--primary);
    background: #374151;
}

.directory-input label {
    color: white !important;
}

.analyze-btn {
    background: linear-gradient(135deg, var(--primary), var(--secondary)) !important;
    color: white !important;
    border: none !important;
    padding: 15px 40px !important;
    border-radius: 50px !important;
    font-size: 1.1em !important;
    font-weight: 600 !important;
    transition: all 0.3s ease !important;
    box-shadow: 0 5px 15px rgba(67, 97, 238, 0.4) !important;
}

.analyze-btn:hover {
    transform: translateY(-2px);
    box-shadow: 0 8px 25px rgba(67, 97, 238, 0.6) !important;
}

.results-section {
    background: var(--dark-card);
    padding: 30px;
    border-radius: 15px;
    box-shadow: 0 8px 16px rgba(0,0,0,0.2);
    margin: 20px 0;
    border: 1px solid var(--dark-border);
}

.results-section h2 {
    color: white;
    margin-bottom: 20px;
    font-weight: 600;
    border-bottom: 3px solid var(--success);
    padding-bottom: 10px;
}

.file-type-badge {
    display: inline-block;
    background: var(--accent);
    color: white;
    padding: 5px 15px;
    border-radius: 20px;
    font-size: 0.9em;
    font-weight: 600;
    margin: 5px;
}

.processing-status {
    background: linear-gradient(135deg, var(--success), #4895ef);
    color: white;
    padding: 20px;
    border-radius: 10px;
    text-align: center;
    margin: 20px 0;
    animation: pulse 2s infinite;
}

@keyframes pulse {
    0% { transform: scale(1); }
    50% { transform: scale(1.02); }
    100% { transform: scale(1); }
}

.footer {
    text-align: center;
    margin-top: 40px;
    padding: 20px;
    color: #a0aec0;
    border-top: 1px solid var(--dark-border);
}

.supported-formats {
    display: flex;
    justify-content: center;
    flex-wrap: wrap;
    gap: 10px;
    margin: 20px 0;
}

.format-badge {
    background: #2d3748;
    padding: 8px 16px;
    border-radius: 20px;
    font-weight: 600;
    color: var(--primary);
    border: 2px solid var(--primary);
}

/* Text color fixes */
label {
    color: white !important;
}

.gr-textbox, .gr-input {
    color: white !important;
}

.gr-markdown {
    color: var(--dark-text) !important;
}

.gr-markdown h1, .gr-markdown h2, .gr-markdown h3 {
    color: white !important;
}
"""

# Create the enhanced UI with DARK theme
with gr.Blocks(css=custom_css, theme=gr.themes.Soft(primary_hue="blue")) as iface:
    with gr.Column(elem_classes="main-container"):
        # Header Section
        with gr.Column(elem_classes="header"):
            gr.Markdown("""
            # 🔒 Advanced Document Redaction System
            ### Secure • Intelligent • Automated
            """)
        
        # Features Grid
        with gr.Row(elem_classes="feature-grid"):
            with gr.Column(elem_classes="feature-card"):
                gr.Markdown("""
                <div class='feature-icon'>🛡️</div>
                <h3>Smart Detection</h3>
                <p>Advanced pattern recognition for sensitive data including PII, credentials, and confidential information</p>
                """)
            
            with gr.Column(elem_classes="feature-card"):
                gr.Markdown("""
                <div class='feature-icon'>🤖</div>
                <h3>AI-Powered Analysis</h3>
                <p>Gemini AI integration for intelligent content analysis and insights generation</p>
                """)
            
            with gr.Column(elem_classes="feature-card"):
                gr.Markdown("""
                <div class='feature-icon'>🎯</div>
                <h3>Multi-Format Support</h3>
                <p>Process PDFs, Word documents, Excel files, PowerPoint, and images seamlessly</p>
                """)
            
            with gr.Column(elem_classes="feature-card"):
                gr.Markdown("""
                <div class='feature-icon'>⚡</div>
                <h3>Batch Processing</h3>
                <p>Automatically analyze and redact multiple files in directory with single click</p>
                """)
        
        # Supported Formats
        with gr.Row():
            with gr.Column():
                gr.Markdown("### 📁 Supported File Formats")
                with gr.Row(elem_classes="supported-formats"):
                    gr.Markdown("<div class='format-badge'>PDF</div>")
                    gr.Markdown("<div class='format-badge'>DOCX</div>")
                    gr.Markdown("<div class='format-badge'>XLSX</div>")
                    gr.Markdown("<div class='format-badge'>PPTX</div>")
                    gr.Markdown("<div class='format-badge'>JPG/PNG</div>")
        
        # Input Section
        with gr.Column(elem_classes="input-section"):
            gr.Markdown("""
            ## 🚀 Get Started
            Enter the directory path containing your documents to begin automated analysis and redaction
            """)
            
            with gr.Row():
                directory_input = gr.Textbox(
                    label="📂 Directory Path",
                    placeholder="C:/Users/YourName/Documents/FilesToAnalyze",
                    lines=1,
                    info="Enter the full path to the directory containing your files",
                    elem_classes="directory-input"
                )
            
            with gr.Row():
                analyze_btn = gr.Button(
                    "🔍 Analyze & Redact All Files", 
                    variant="primary", 
                    size="lg",
                    elem_classes="analyze-btn"
                )
        
        # Results Section
        with gr.Column(elem_classes="results-section"):
            gr.Markdown("## 📊 Analysis Results")
            
            with gr.Row():
                output_html = gr.HTML(
                    label="📋 Document Analysis Report",
                    show_label=True
                )
        
        # Footer
        with gr.Column(elem_classes="footer"):
            gr.Markdown("""
            ---
            **🔐 Security First** • **🤖 Powered by AI** • **⚡ Built for Performance**
            
            *Your documents are processed securely and never stored on our servers*
            """)
    
    # Connect the functionality
    analyze_btn.click(
        fn=process_directory,
        inputs=[directory_input],
        outputs=[output_html]
    )

if __name__ == "__main__":
    iface.launch(
        share=False, 
        server_name="0.0.0.0", 
        server_port=7860,
        show_error=True
    )