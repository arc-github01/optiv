
import io
import fitz
import docx
import pandas as pd
import openpyxl
from openpyxl.styles import PatternFill, Font
from pptx import Presentation
from PIL import Image, ImageFilter, ImageDraw, ImageStat
import pytesseract
import cv2
import numpy as np
from utils.text_utils import is_sensitive_text, redact_text
from utils.image_utils import (
    blur_faces, create_clean_image
)


def has_text_content(image, threshold=0.3):
    """
    Detect if image contains text/logos using multiple methods
    Returns True if text is detected
    """
    try:
        # Convert PIL to numpy array for OpenCV
        img_array = np.array(image)
        
        # Method 1: OCR detection
        text = pytesseract.image_to_string(image, config='--psm 6')
        if len(text.strip()) > 2:  # If any text is detected
            return True
        
        # Method 2: Edge detection (logos have high edge density)
        gray = cv2.cvtColor(img_array, cv2.COLOR_RGB2GRAY)
        edges = cv2.Canny(gray, 50, 150)
        edge_percentage = np.sum(edges > 0) / edges.size
        
        if edge_percentage > threshold:
            return True
        
        # Method 3: Color variance (logos often have distinct colors)
        stat = ImageStat.Stat(image)
        variance = sum(stat.var) / len(stat.var)
        
        if variance > 1000:  # High variance suggests logo/text
            return True
            
        return False
        
    except Exception as e:
        print(f"Error in text detection: {e}")
        return False


def is_logo_like(image):
    """
    Determine if an image is likely a logo based on characteristics
    """
    try:
        # Check aspect ratio (logos often wider than tall)
        width, height = image.size
        aspect_ratio = width / height if height > 0 else 1
        
        # Check if image is small (logos are usually small)
        total_pixels = width * height
        
        # Check transparency (logos often have transparency)
        has_alpha = image.mode == 'RGBA'
        
        # Check for text
        has_text = has_text_content(image, threshold=0.2)
        
        # Logo characteristics:
        # - Small to medium size
        # - Often has text
        # - May have specific aspect ratios
        is_small = total_pixels < 500000  # Less than 500k pixels
        
        if has_text and is_small:
            return True
        
        if has_text and (aspect_ratio > 2 or aspect_ratio < 0.5):
            return True
            
        return False
        
    except Exception as e:
        print(f"Error in logo detection: {e}")
        return False


def black_out_image(image):
    """
    Completely black out an image
    """
    return Image.new('RGB', image.size, (0, 0, 0))


def heavy_blur_image(image, radius=50):
    """
    Apply very heavy blur that makes content unrecognizable
    """
    # Apply multiple passes of blur
    blurred = image
    for _ in range(5):
        blurred = blurred.filter(ImageFilter.GaussianBlur(radius))
    return blurred


def detect_and_redact_logos(image):
    """
    Detect and completely redact logos and text in images
    """
    try:
        # Check if image contains text/logo content
        if is_logo_like(image):
            print("Logo/text detected - applying complete blackout")
            # For logos/text, black out completely
            return black_out_image(image)
        
        # Check for any text content
        if has_text_content(image, threshold=0.15):
            print("Text content detected - applying heavy blur")
            # Apply very heavy blur
            return heavy_blur_image(image, radius=50)
        
        # If no text detected, return original
        return image
        
    except Exception as e:
        print(f"Error in logo redaction: {e}")
        return image


def blur_small_images(image, max_dimension=200):
    """
    Heavily blur or black out small images (likely logos)
    """
    try:
        width, height = image.size
        
        # If image is small, it's likely a logo - black it out
        if width < max_dimension and height < max_dimension:
            print(f"Small image detected ({width}x{height}) - treating as logo")
            return black_out_image(image)
        
        # If image is small in one dimension, check for text
        if width < max_dimension or height < max_dimension:
            print(f"Small dimension detected - checking for logo/text")
            if has_text_content(image, threshold=0.1):
                return black_out_image(image)
            return heavy_blur_image(image, radius=40)
        
        return image
        
    except Exception as e:
        print(f"Error processing small images: {e}")
        return image


def detect_and_blur_codes(image):
    """
    Detect QR codes, barcodes, and apply heavy redaction
    """
    try:
        # First check for logos/text and redact them
        image = detect_and_redact_logos(image)
        
        # Then check for QR codes and barcodes
        img_array = np.array(image)
        detector = cv2.QRCodeDetector()
        data, bbox, _ = detector.detectAndDecode(img_array)
        
        if bbox is not None:
            pil_image = image.copy()
            draw = ImageDraw.Draw(pil_image)
            points = bbox[0].astype(int)
            
            x_min = int(points[:, 0].min())
            y_min = int(points[:, 1].min())
            x_max = int(points[:, 0].max())
            y_max = int(points[:, 1].max())
            draw.rectangle([x_min, y_min, x_max, y_max], fill=(0, 0, 0))
            
            return pil_image
        
        return image
    except Exception as e:
        print(f"Error in code detection: {e}")
        return image


def redact_pdf(input_path, output_path):
    """Create redacted PDF by overlaying black rectangles on sensitive text"""
    try:
        doc = fitz.open(input_path)
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            
            # Redact text - PROPERLY removes text, not just covers it
            text_instances = page.get_text("dict")
            blocks = text_instances.get("blocks", [])
            
            # Collect all areas to redact
            redaction_areas = []
            
            for block in blocks:
                if "lines" in block:
                    for line in block["lines"]:
                        line_text = " ".join([span.get("text", "") for span in line["spans"]])
                        if is_sensitive_text(line_text):
                            # Redact entire line
                            for span in line["spans"]:
                                rect = fitz.Rect(span["bbox"])
                                redaction_areas.append(rect)
                        else:
                            # Check individual spans
                            for span in line["spans"]:
                                text = span.get("text", "")
                                if is_sensitive_text(text):
                                    rect = fitz.Rect(span["bbox"])
                                    redaction_areas.append(rect)
            
            # Apply proper redaction annotations (this actually removes the text)
            for rect in redaction_areas:
                # Add redaction annotation
                annot = page.add_redact_annot(rect, fill=(0, 0, 0))
            
            # Apply all redactions (this permanently removes the text)
            page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_NONE)
            
            # Process images - CRITICAL: Delete original images first
            image_list = page.get_images(full=True)
            
            # Store image info before deletion
            images_to_process = []
            for img_index, img_info in enumerate(image_list):
                try:
                    xref = img_info[0]
                    
                    # Get image rectangle BEFORE deletion
                    try:
                        img_rect = page.get_image_rects(img_info)[0]
                    except (IndexError, Exception):
                        try:
                            img_rect = page.get_image_bbox(img_info)
                        except:
                            print(f"Skipping image {img_index} - cannot get rectangle")
                            continue
                    
                    if not (img_rect and img_rect.is_valid and not img_rect.is_empty):
                        print(f"Invalid image rectangle for image {img_index}")
                        continue
                    
                    # Extract image data
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    
                    images_to_process.append({
                        'xref': xref,
                        'rect': img_rect,
                        'bytes': image_bytes
                    })
                    
                except Exception as img_error:
                    print(f"Error preparing image {img_index} on page {page_num}: {img_error}")
                    continue
            
            # Delete all original images from the page
            for img_data in images_to_process:
                try:
                    page.delete_image(img_data['xref'])
                except Exception as e:
                    print(f"Error deleting image: {e}")
            
            # Now process and insert the cleaned images
            for img_data in images_to_process:
                try:
                    pil_image = Image.open(io.BytesIO(img_data['bytes']))
                    processed_image = pil_image.convert("RGB")
                    
                    # Apply all processing in the correct order
                    # 1. Check and black out small images (likely logos)
                    processed_image = blur_small_images(processed_image)
                    
                    # 2. Detect and redact any logos/text (heavy blur or blackout)
                    processed_image = detect_and_blur_codes(processed_image)
                    
                    # 3. Blur faces
                    processed_image = blur_faces(processed_image)
                    
                    # 4. Create clean version without metadata
                    processed_image = create_clean_image(processed_image)
                    
                    # Save processed image
                    img_byte_arr = io.BytesIO()
                    processed_image.save(img_byte_arr, format='PNG', exif=b'')
                    img_byte_arr.seek(0)
                    
                    # Insert the processed image at the original location
                    page.insert_image(img_data['rect'], stream=img_byte_arr.getvalue(), keep_proportion=False)
                    
                except Exception as img_error:
                    print(f"Error processing/inserting image on page {page_num}: {img_error}")
                    continue
        
        doc.save(output_path)
        doc.close()
        return output_path
        
    except Exception as e:
        print(f"Error redacting PDF: {e}")
        import traceback
        traceback.print_exc()
        return None


def redact_docx(input_path, output_path):
    """Create redacted DOCX by replacing sensitive text"""
    try:
        doc = docx.Document(input_path)
        
        for paragraph in doc.paragraphs:
            if paragraph.text:
                redacted_text = redact_text(paragraph.text)
                if redacted_text != paragraph.text:
                    paragraph.text = redacted_text
        
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text:
                        redacted_text = redact_text(cell.text)
                        if redacted_text != cell.text:
                            cell.text = redacted_text
        
        doc.save(output_path)
        return output_path
    except Exception as e:
        print(f"Error redacting DOCX: {e}")
        import traceback
        traceback.print_exc()
        return None


def redact_excel(input_path, output_path):
    """Redact sensitive information in Excel files"""
    try:
        if input_path.endswith('.xlsx'):
            wb = openpyxl.load_workbook(input_path)
            black_fill = PatternFill(start_color="000000", end_color="000000", fill_type="solid")
            white_font = Font(color="FFFFFF")
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value is None:
                            continue
                        
                        if is_sensitive_text(str(cell.value)):
                            cell.fill = black_fill
                            cell.font = white_font
                            cell.value = "[REDACTED]"
            
            wb.save(output_path)
            return output_path
        else:
            output_path = output_path.replace('.xls', '.xlsx')
            excel_file = pd.ExcelFile(input_path)
            with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
                for sheet_name in excel_file.sheet_names:
                    df = pd.read_excel(input_path, sheet_name=sheet_name)
                    
                    for col in df.columns:
                        df[col] = df[col].apply(
                            lambda x: "[REDACTED]" if pd.notna(x) and is_sensitive_text(str(x)) else x
                        )
                    
                    df.to_excel(writer, sheet_name=sheet_name, index=False)
            
            return output_path
            
    except Exception as e:
        print(f"Error redacting Excel: {e}")
        import traceback
        traceback.print_exc()
        return None


def redact_ppt(input_path, output_path):
    """Redact sensitive information in PowerPoint files"""
    try:
        prs = Presentation(input_path)
        
        for slide in prs.slides:
            shapes_to_replace = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    original_text = shape.text
                    redacted_text = redact_text(original_text)
                    
                    if shape.has_text_frame and original_text != redacted_text:
                        shape.text_frame.clear()
                        p = shape.text_frame.paragraphs[0] if shape.text_frame.paragraphs else shape.text_frame.add_paragraph()
                        run = p.add_run()
                        run.text = redacted_text
                
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            if cell.text:
                                original_text = cell.text
                                redacted_text = redact_text(original_text)
                                if original_text != redacted_text:
                                    cell.text_frame.clear()
                                    p = cell.text_frame.paragraphs[0] if cell.text_frame.paragraphs else cell.text_frame.add_paragraph()
                                    run = p.add_run()
                                    run.text = redacted_text
                
                if shape.shape_type == 13:  # Picture shape
                    shapes_to_replace.append({
                        'shape': shape,
                        'left': shape.left,
                        'top': shape.top,
                        'width': shape.width,
                        'height': shape.height
                    })

            # Process images - CRITICAL: Add new image BEFORE removing old one
            for shape_info in shapes_to_replace:
                try:
                    shape = shape_info['shape']
                    image = shape.image
                    image_bytes = image.blob
                    pil_image = Image.open(io.BytesIO(image_bytes))
                    processed_image = pil_image.convert("RGB")
                    
                    # Apply all processing in the correct order
                    # 1. Check and black out small images (likely logos)
                    processed_image = blur_small_images(processed_image)
                    
                    # 2. Detect and redact any logos/text (heavy blur or blackout)
                    processed_image = detect_and_blur_codes(processed_image)
                    
                    # 3. Blur faces
                    processed_image = blur_faces(processed_image)
                    
                    # 4. Create clean version without metadata
                    processed_image = create_clean_image(processed_image)
                    
                    with io.BytesIO() as output:
                        processed_image.save(output, format='PNG', exif=b'')
                        output.seek(0)
                        
                        # Add the new picture first
                        new_pic = slide.shapes.add_picture(
                            output, 
                            shape_info['left'], 
                            shape_info['top'], 
                            width=shape_info['width'], 
                            height=shape_info['height']
                        )
                        
                        # Now remove the old shape
                        sp = shape._element
                        sp.getparent().remove(sp)
                        
                except Exception as img_error:
                    print(f"Error processing image in slide: {img_error}")
                    import traceback
                    traceback.print_exc()
                    continue

        prs.save(output_path)
        return output_path
    except Exception as e:
        print(f"Error redacting PowerPoint: {e}")
        import traceback
        traceback.print_exc()
        return None