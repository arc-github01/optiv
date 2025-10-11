"""
Document content extraction utilities with secure PII redaction
"""
import fitz
import docx
import pandas as pd
import openpyxl
from pptx import Presentation
from utils.image_utils import get_safe_text_for_llm, get_safe_image_for_llm, process_image_secure
from utils.text_utils import redact_text


def extract_from_pdf(path):
    """Extract text from PDF file with redaction"""
    try:
        doc = fitz.open(path)
        text = ""
        for page in doc:
            text += page.get_text()
        doc.close()
        
        # Apply text redaction to extracted PDF text
        text = redact_text(text, full_text=text)
        return text
    except Exception as e:
        print(f"Error reading PDF: {e}")
        import traceback
        traceback.print_exc()
        return ""


def extract_from_docx(path):
    """Extract text from DOCX file with redaction"""
    try:
        docx_file = docx.Document(path)
        text_content = []
        
        for paragraph in docx_file.paragraphs:
            text_content.append(paragraph.text)
        
        for table in docx_file.tables:
            for row in table.rows:
                row_text = [cell.text for cell in row.cells]
                text_content.append("\t".join(row_text))
        
        full_text = "\n".join(text_content)
        
        # Apply text redaction to extracted DOCX text
        full_text = redact_text(full_text, full_text=full_text)
        return full_text
    except Exception as e:
        print(f"Error reading DOCX: {e}")
        import traceback
        traceback.print_exc()
        return ""


def extract_from_excel(path):
    """Extract text from Excel file with redaction"""
    try:
        if path.endswith('.xlsx'):
            wb = openpyxl.load_workbook(path, data_only=True)
            text_content = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_content.append(f"\n=== Sheet: {sheet_name} ===\n")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = [str(cell) if cell is not None else "" for cell in row]
                    if any(row_text):
                        text_content.append("\t".join(row_text))
            
            full_text = "\n".join(text_content)
        else:
            excel_file = pd.ExcelFile(path)
            text_content = []
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(path, sheet_name=sheet_name)
                text_content.append(f"\n=== Sheet: {sheet_name} ===\n")
                df = df.fillna('')
                text_content.append(df.to_string(index=False))
            full_text = "\n".join(text_content)
        
        # Apply text redaction to extracted Excel text
        full_text = redact_text(full_text, full_text=full_text)
        return full_text
    except Exception as e:
        print(f"Error reading Excel file: {e}")
        import traceback
        traceback.print_exc()
        return f"Error reading Excel file: {e}"


def extract_from_ppt(path):
    """Extract text from PowerPoint file with redaction"""
    try:
        prs = Presentation(path)
        text_content = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            text_content.append(f"\n=== Slide {slide_num} ===\n")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_content.append(shape.text)
                
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = [cell.text for cell in row.cells]
                        text_content.append("\t".join(row_text))
                
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text:
                                text_content.append(run.text)
        
        full_text = "\n".join(text_content)
        
        # Apply text redaction to extracted PowerPoint text
        full_text = redact_text(full_text, full_text=full_text)
        return full_text
    except Exception as e:
        print(f"Error reading PowerPoint file: {e}")
        import traceback
        traceback.print_exc()
        return f"Error reading PowerPoint file: {e}"


def extract_from_image(path):
    """
    Extract text from image file with SECURE pre-redaction.
    
    This function ensures sensitive information is redacted BEFORE
    any text is extracted and exposed.
    
    Args:
        path: Path to the image file
    
    Returns:
        str: Securely redacted text extracted from the image
    """
    try:
        # Use secure extraction - redacts image FIRST, then extracts text
        text = get_safe_text_for_llm(path)
        return text
    except Exception as e:
        print(f"Error extracting text from image: {e}")
        import traceback
        traceback.print_exc()
        return ""


def extract_from_image_with_validation(path):
    """
    Extract text from image with full validation reporting.
    
    Use this when you need to verify redaction effectiveness.
    
    Args:
        path: Path to the image file
    
    Returns:
        dict: Contains 'text', 'image', 'validation', and 'is_valid'
    """
    try:
        result = process_image_secure(path, return_original=False)
        
        is_valid, issues = result['validation']
        
        return {
            'text': result['redacted_text'],
            'image': result['redacted_image'],
            'validation': result['validation'],
            'is_valid': is_valid,
            'issues': issues
        }
    except Exception as e:
        print(f"Error extracting text from image: {e}")
        import traceback
        traceback.print_exc()
        return {
            'text': "",
            'image': None,
            'validation': (False, [str(e)]),
            'is_valid': False,
            'issues': [str(e)]
        }


# Backward compatibility - but logs warning
def extract_text_from_image(path):
    """
    DEPRECATED: Use extract_from_image() instead.
    
    This function is kept for backward compatibility but redirects
    to the secure version.
    """
    print("⚠️  WARNING: extract_text_from_image() is deprecated.")
    print("   Use extract_from_image() or extract_from_image_with_validation() instead.")
    return extract_from_image(path)